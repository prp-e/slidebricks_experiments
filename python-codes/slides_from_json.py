import json
from pptx import Presentation

def create_presentation_from_json(json_data):
    # Create a new presentation object
    presentation = Presentation()

    # Add the title slide
    title_slide_layout = presentation.slide_layouts[0]  # Title Slide layout
    slide = presentation.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = json_data["title"]
    subtitle.text = json_data["subtitle"]

    # Iterate over each slide in the JSON
    for slide_data in json_data["slides"]:
        slide_layout = None
        slide = None

        if slide_data["type"] == "title_slide":
            slide_layout = presentation.slide_layouts[0]  # Title Slide layout
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = slide_data["title"]
            subtitle.text = slide_data["subtitle"]

        elif slide_data["type"] == "title_and_content":
            slide_layout = presentation.slide_layouts[1]  # Title and Content layout
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = slide_data["title"]
            content.text = "\n".join(slide_data["content"])

        elif slide_data["type"] == "section_header":
            slide_layout = presentation.slide_layouts[2]  # Section Header layout
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = slide_data["title"]
            subtitle.text = slide_data["subtitle"]

        elif slide_data["type"] == "two_content":
            slide_layout = presentation.slide_layouts[3]  # Two Content layout
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            left_content = slide.placeholders[1]
            right_content = slide.placeholders[2]
            title.text = slide_data["title"]
            left_content.text = "\n".join(slide_data["left_content"])
            right_content.text = "\n".join(slide_data["right_content"])

        elif slide_data["type"] == "comparison":
            slide_layout = presentation.slide_layouts[4]  # Comparison layout
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            left_title = slide.placeholders[1]
            left_content = slide.placeholders[2]
            right_title = slide.placeholders[3]
            right_content = slide.placeholders[4]
            title.text = slide_data["title"]
            left_title.text = slide_data["left_title"]
            left_content.text = "\n".join(slide_data["left_content"])
            right_title.text = slide_data["right_title"]
            right_content.text = "\n".join(slide_data["right_content"])

        elif slide_data["type"] == "title_only":
            slide_layout = presentation.slide_layouts[5]  # Title Only layout
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = slide_data["title"]

        elif slide_data["type"] == "blank":
            slide_layout = presentation.slide_layouts[6]  # Blank layout
            slide = presentation.slides.add_slide(slide_layout)

        elif slide_data["type"] == "content_with_caption":
            slide_layout = presentation.slide_layouts[7]  # Content with Caption layout
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            caption = slide.placeholders[2]
            title.text = slide_data["title"]
            content.text = slide_data["content"]
            caption.text = slide_data["caption"]

        elif slide_data["type"] == "picture_with_caption":
            slide_layout = presentation.slide_layouts[8]  # Picture with Caption layout
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            caption = slide.placeholders[1]
            title.text = slide_data["title"]
            caption.text = slide_data["caption"]

    # Add the end page
    end_slide_layout = presentation.slide_layouts[1]  # Title and Content layout
    end_slide = presentation.slides.add_slide(end_slide_layout)
    title = end_slide.shapes.title
    content = end_slide.placeholders[1]
    title.text = "Conclusion"
    content.text = json_data["endpage"]

    # Save the presentation
    presentation.save("comprehensive_presentation.pptx")
    print("Presentation created successfully!")

# Example usage
json_input = '''
{
    "title": "Comprehensive Presentation",
    "subtitle": "Exploring All Slide Types",
    "slides": [
        {
            "type": "title_slide",
            "title": "Introduction",
            "subtitle": "Overview of Slide Types"
        },
        {
            "type": "title_and_content",
            "title": "Agenda",
            "content": [
                "Introduction to Slide Types",
                "Detailed Examples",
                "Summary"
            ]
        },
        {
            "type": "section_header",
            "title": "Main Section",
            "subtitle": "Detailed Exploration"
        },
        {
            "type": "two_content",
            "title": "Comparison",
            "left_content": [
                "Point A1",
                "Point A2"
            ],
            "right_content": [
                "Point B1",
                "Point B2"
            ]
        },
        {
            "type": "comparison",
            "title": "Feature Comparison",
            "left_title": "Feature A",
            "left_content": [
                "Detail A1",
                "Detail A2"
            ],
            "right_title": "Feature B",
            "right_content": [
                "Detail B1",
                "Detail B2"
            ]
        },
        {
            "type": "title_only",
            "title": "Summary"
        },
        {
            "type": "blank"
        },
        {
            "type": "content_with_caption",
            "title": "Content with Caption",
            "content": "This is a content area with a caption.",
            "caption": "Caption text here."
        },
        {
            "type": "picture_with_caption",
            "title": "Picture with Caption",
            "caption": "This is a caption for the picture."
        }
    ],
    "endpage": "Thank you for your attention!"
}
'''

# Load the JSON data
json_data = json.loads(json_input)

# Create the presentation
create_presentation_from_json(json_data)