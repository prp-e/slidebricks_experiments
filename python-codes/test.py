from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
presentation = Presentation()

# Slide 1: Title Slide
slide_1 = presentation.slides.add_slide(presentation.slide_layouts[0])  # Title Slide layout
title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1]

title_1.text = "Welcome to My Presentation"
subtitle_1.text = "An Overview of Python-PPTX"

# Slide 2: Introduction
slide_2 = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
title_2 = slide_2.shapes.title
content_2 = slide_2.placeholders[1]

title_2.text = "Introduction"
content_2.text = "In this presentation, we will cover:\n- Basics of python-pptx\n- Creating slides\n- Adding content\n- Saving presentations"

# Slide 3: Features
slide_3 = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
title_3 = slide_3.shapes.title
content_3 = slide_3.placeholders[1]

title_3.text = "Features of python-pptx"
content_3.text = "Some features include:\n- Easy to use API\n- Supports images and charts\n- Customizable layouts\n- Open and save presentations"

# Slide 4: Example Code
slide_4 = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
title_4 = slide_4.shapes.title
content_4 = slide_4.placeholders[1]

title_4.text = "Example Code"
content_4.text = "Here's a simple example:\n\n```python\nfrom pptx import Presentation\n# Create a presentation\npresentation = Presentation()\n```"

# Slide 5: Conclusion
slide_5 = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
title_5 = slide_5.shapes.title
content_5 = slide_5.placeholders[1]

title_5.text = "Conclusion"
content_5.text = "Thank you for your attention!\n- Questions?\n- Feedback?"

# Save the presentation
presentation.save("simple_presentation.pptx")

print("Presentation created successfully!")
